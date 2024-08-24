'''
Author: J , jwsun1987@gmail.com
Date: 2024-02-09 17:48:56
LastEditors: J , jwsun1987@gmail.com
Description: 
Copyright: Copyright (c) 2024 by jwsun1987@gmail.com. All Rights Reserved.
'''


import pandas as pd

try:
    from pptx.chart.data import ChartData, XyChartData
    from pptx import enum as pptx_enum
except ImportError:
    import sys

    sys.path.append(
        r"\\merlin\lib_isg\\28.Alternative Data\Code\python-quant\libs\common\reporting\slides\python-pptx-master"
    )
    from pptx.chart.data import ChartData, XyChartData
    from pptx import enum as pptx_enum

from common.data import *
from common.dates import *
from utils.tools import *
from .lo_slides_style import *
from common.reporting.style.lo_color import *
from common.reporting.style.lo_style import *


from .chart_style import *
from .table_style import *

LO_DEFAULT_FONT_COLOR = RGBColor(57, 46, 44)

def plot_textbox(shape, pos_x, pos_y, width, height, params={}):
    print('text box params')
    print(params)
    # add a text box
    text_box = shape.add_textbox(pos_x, pos_y, width, height)
    # set text boarder line color
    if params.get("line", False):
        text_box.line.color.rgb = params.get("line_color", LO_DEFAULT_FONT_COLOR)
        text_box.line.width = params.get("line_width", TEXTBOX_LINE_WIDTH)

    # set the text formart
    texts = to_list(params.get("text", [""]))
    for idx, text in enumerate(texts):
        if idx > 0:
            text_box.text_frame.add_paragraph()
        text_box.text_frame.paragraphs[idx].text = text
        text_box.text_frame.paragraphs[idx].alignment = params.get(
            "align", pptx_enum.text.PP_ALIGN.LEFT
        )
        text_box.text_frame.paragraphs[idx].font.name = params.get("font", FONT_NAME)
        text_box.text_frame.paragraphs[idx].font.size = params.get(
            "font_size", CHART_INFO_FONT_SIZE
        )
        text_box.text_frame.paragraphs[idx].font.bold = params.get("font_bold", True)
        text_box.text_frame.paragraphs[idx].font.color.rgb = params.get(
            "font_color", RGBColor(57, 46, 44)# LO_DEFAULT_FONT_COLOR
        )
    if params.get("solid", False):
        text_box.fill.solid()
        text_box.fill.fore_color.rgb = params.get("color", LO_LIGHTTAN)


def plot_chart(
    shape, chart_data, pos_x, pos_y, width, height, chart_type_str: str, params={}
):
    #print(params)
    # add caption
    if params.get("caption", {}):
        plot_textbox(
            shape,
            pos_x,
            pos_y,
            params.get("caption", {}).get("width", width),
            params.get("caption", {}).get("height", CAPTION_HEIGHT),
            params.get("caption", {}),
        )
        pos_y += params.get("caption_height", CAPTION_HEIGHT) * 0.75
        height -= params.get("caption_height", CAPTION_HEIGHT) * 0.75

    if params.get("pos_y", None) is not None:
        pos_y = params.get("pos_y")
    if params.get("height", None) is not None:
        height = params.get("height")

    # translate chart_type to chart.XL_CHART_TYPE
    chart_type = translate_chart_type(chart_type_str)

    # plot the chart
    chart = shape.add_chart(
        chart_type,
        pos_x,
        pos_y,
        width,
        height,
        chart_data,
    ).chart

    # set font
    set_chart_font(chart, params.get("style", {}))

    # set category axis (x)
    set_chart_axis(
        chart.category_axis,
        params.get("category_axis", {}),
    )
    set_chart_axis_ticks(
        chart.category_axis,
        params.get("category_axis", {}),
    )

    # set value axis (y)
    value_axis_params = add_params(
        params.get(
            "value_axis", {"major_tick_mark": pptx_enum.chart.XL_TICK_MARK.OUTSIDE}
        ),
        {"major_tick_mark": pptx_enum.chart.XL_TICK_MARK.OUTSIDE},
    )
    set_chart_axis(chart.value_axis, value_axis_params)
    set_chart_axis_ticks(chart.value_axis, value_axis_params)

    # set data label
    if chart_type_str != "scatter":
        set_chart_data_labels(chart, params.get("data_labels", {}))

    # set legend
    set_chart_legend(chart, params.get("legend", {}))

    # set title
    set_chart_title(chart, params.get("title", {}))

    # set series plot params
    set_chart_style(chart, chart_type, params.get("style", {}))


def plot_clustered_columns(
    slide, data: pd.DataFrame, pos_x, pos_y, width, height, params={}
):
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = data.index
    for col in data.columns:
        chart_data.add_series(col, data[col].values)

    # add chart to slide --------------------
    plot_chart(slide.shapes, chart_data, pos_x, pos_y, width, height, "column", params)


def plot_scatter_chart(
    slide, data_x, data_y, pos_x, pos_y, width, height, chart_params={}
):
    chart_data = XyChartData()

    cd = chart_data.add_series(
        chart_params.get("title", {}).get("text", ""),
        number_format=chart_params.get("data_labels", {}).get("number_format", None),
    )
    for x, y in list(zip(data_x.values, data_y.values)):
        cd.add_data_point(
            x,
            y,
            number_format=chart_params.get("data_labels", {}).get(
                "number_format", None
            ),
        )
    # chart_data.x_values = data_x.values
    # chart_data.y_values = data_y.values
    plot_chart(
        slide.shapes, chart_data, pos_x, pos_y, width, height, "scatter", chart_params
    )


def plot_table(slide, df: pd.DataFrame, pos_x, pos_y, width, height, params={}):
    n_header_lvl = df.columns.nlevels
    n_index_lvl = df.index.nlevels
    # add a new table
    table_shape = slide.shapes.add_table(
        (df.shape[0] + n_header_lvl),
        (df.shape[1] + n_index_lvl),
        int(pos_x),
        int(pos_y),
        int(width),
        int(height),
    )
    # set table style
    set_table_style(table_shape, params)

    # fill table contents
    table = table_shape.table
    set_table_content(table, df, params)
