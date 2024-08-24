'''
Author: J , jwsun1987@gmail.com
Date: 2024-02-09 17:45:12
LastEditors: J , jwsun1987@gmail.com
Description: 
Copyright: Copyright (c) 2024 by jwsun1987@gmail.com. All Rights Reserved.
'''


import datetime

from typing import List, Union, Dict, Optional, Sequence
try:
    from pptx.chart.data import ChartData, XyChartData
except ImportError:
    import sys

    sys.path.append(
        r"\\merlin\lib_isg\\28.Alternative Data\Code\python-quant\libs\common\reporting\slides\python-pptx-master"
    )
    from pptx.chart.data import ChartData, XyChartData

from common.dates import *
from utils.tools import *

from .lo_slides_style import *
from .chart_data import *
from .pptx_wrapper import *

ROUND_BASE = 5


def add_footnotes(
    slide,
    footnotes: list,
    pos_x=BODY_X,
    pos_y=BODY_Y + BODY_HEIGHT + MARGIN_VER,
    width=BODY_WIDTH,
    height=2 * MARGIN_VER,
    params: dict = {},
):
    if not footnotes:
        footnotes = ["Source: LombardOdier Computation"]

    plot_textbox(
        slide.shapes,
        pos_x,
        pos_y,
        width,
        height,
        add_params(
            params,
            {
                "text": footnotes,
                "align": PP_ALIGN.LEFT,
                "font_size": FOOTNOTE_FONT_SIZE,
                "font_bold": False,
            },
        ),
    )


def add_new_content_page(presentation, params: dict = {}):
    new_slide = presentation.slides.add_slide(
        presentation.slide_masters[0].slide_layouts.get_by_name("Contents")
    )

    for s in new_slide.shapes:
        if "Title" in s.name:
            texts = to_list(params.get("title", {}).get("text", ""))
            for idx, text in enumerate(texts):
                if idx > 0:
                    s.text_frame.add_paragraph()
                    s.text_frame.paragraphs[idx].font.size = Pt(18)
                    s.text_frame.paragraphs[idx].font.bold = False
                else:
                    s.text_frame.paragraphs[idx].font.bold = params.get(
                        "title", {}
                    ).get("bold", True)
                s.text_frame.paragraphs[idx].font.name = params.get("title", {}).get(
                    "font", FONT_NAME
                )
                s.text_frame.paragraphs[idx].text = text
    return new_slide


def update_title_page(slide, title: str, authors: str, date=None, subtitle=None):
    for shape in slide.shapes:
        if "Title" in shape.name:
            shape.text = title
        elif "Subtitle" in shape.name:
            shape.text = "" if subtitle is None else subtitle
        elif "Date" in shape.name:
            date = datetime.datetime.today() if date is None else pd.to_datetime(date)
            shape.text = (
                str(date.date().day).zfill(2)
                + "."
                + str(date.date().month).zfill(2)
                + "."
                + str(date.date().year)
            )
        elif "Text" in shape.name:
            shape.text = authors


def append_final_pages(presentation):
    presentation.slides.add_slide(
        presentation.slide_masters[0].slide_layouts.get_by_name("Disclaimer")
    )
    presentation.slides.add_slide(
        presentation.slide_masters[0].slide_layouts.get_by_name("Last Page")
    )


def add_table_page(
    presentation, data_table: Union[pd.DataFrame, pd.Series], params: dict = {}
):
    new_slide = add_new_content_page(presentation, params)

    # set table position and dimension
    table_x = BODY_X
    table_y = TITLE_Y + TITLE_HEIGHT
    table_width = BODY_WIDTH
    table_height = BODY_HEIGHT
    # add the table
    plot_table(
        new_slide, data_table, table_x, table_y, table_width, table_height, params
    )
    add_footnotes(new_slide, params.get("footnotes", []))


def add_text_chart_grid_page(
    presentation, data: dict, rows: int = 2, cols: int = 2, params: dict = {}
):
    # chart grids
    # grid size
    grid_w = BODY_WIDTH / cols
    grid_h = BODY_HEIGHT / rows
    text_box_number = params.get("text_box", 1)

    page_ind = -1
    for ind, d_key in enumerate(data.keys()):
        ind_per_page = ind % (rows * cols - text_box_number)
        if ind_per_page == 0:
            new_slide = add_new_content_page(presentation, params)
            if text_box_number > 0:
                for text_box in range(text_box_number):
                    txt_x_multiple = text_box % cols
                    txt_y_multiple = text_box // cols
                    txt_w = BODY_X + txt_x_multiple * grid_w
                    txt_h = BODY_Y + txt_y_multiple * grid_h
                    txt = new_slide.shapes.add_textbox(txt_w, txt_h, grid_w, grid_h)
                    txt.fill.solid()
                    txt.fill.fore_color.rgb = params.get("txtbox_color", LO_LIGHTTAN)
            page_ind += 1

        # grid position
        x_multiple = (ind_per_page + text_box_number) % cols
        y_multiple = (ind_per_page + text_box_number) // cols
        d_pos_x = BODY_X + x_multiple * grid_w
        d_pos_y = BODY_Y + y_multiple * grid_h

        # prepare chart data
        d_chart = line_chart_data(data.get(d_key), params)
        # prepare chart params
        d_chart_params = add_params(
            params.get("chart", {"title": {"on": False}}),
            {
                **d_chart["params"],
                **{
                    "caption": {"text": d_key, "font_size": CHART_CAPTION_FONT_SIZE},
                    "style": {"line": {"country": False}},
                },
            },
        )
        # plot the chart
        plot_chart(
            new_slide.shapes,
            d_chart["data"],
            d_pos_x,
            d_pos_y,
            grid_w,
            grid_h,
            d_chart_params.get("type", "line"),
            d_chart_params,
        )
    add_footnotes(new_slide, params.get("footnotes", []))


def add_chart_grid_page(
    presentation,
    data: dict,
    pos_x=BODY_X,
    pos_y=BODY_Y,
    width=BODY_WIDTH,
    height=BODY_HEIGHT,
    rows: int = 1,
    cols: int = 2,
    new_page: bool = True,
    params: dict = {},
):
    # chart grids
    # grid size
    grid_w = width / cols
    grid_h = height / rows

    page_ind = -1
    for ind, d_key in enumerate(data.keys()):
        print(d_key)
        if ind % (rows * cols) == 0:
            if ind == 0:
                if new_page:
                    new_slide = add_new_content_page(presentation, params)
                else:
                    new_slide = presentation.slides[-1]
            else:
                new_slide = add_new_content_page(presentation, params)
            page_ind += 1

        # grid position
        x_multiple = ind % cols
        y_multiple = ind // cols
        d_pos_x = pos_x + x_multiple * grid_w
        d_pos_y = pos_y + (y_multiple - page_ind * rows) * grid_h

        # prepare chart data
        d_chart = line_chart_data(data.get(d_key), params)

        # prepare chart params
        caption_params = add_params(
            params.get("chart", {}).get("caption",{}),
            {"text": d_key, "font_size": CHART_CAPTION_FONT_SIZE} 
        )

        d_chart_params = add_params(
            params.get("chart", {"title": {"on": False}}),
            {
                **d_chart["params"],
                **{
                    "caption":caption_params,
                },
                **params.get("chart",{}).get("style", {"line": {"country": False}}),
            },
        )

        # plot the chart
        plot_chart(
            new_slide.shapes,
            d_chart["data"],
            d_pos_x,
            d_pos_y,
            grid_w,
            grid_h,
            d_chart_params.get("type", "line"),
            d_chart_params,
        )
        add_footnotes(new_slide, params.get("footnotes", []))


def add_textbox_grid_page(
    presentation,
    data: dict,
    pos_x=BODY_X,
    pos_y=BODY_Y,
    width=BODY_WIDTH,
    height=BODY_HEIGHT,
    rows: int = 1,
    cols: int = 2,
    new_page: bool = True,
    params: dict = {},
):
    # chart grids
    # grid size
    grid_w = width / cols
    grid_h = height / rows

    page_ind = -1
    for ind, d_key in enumerate(data.keys()):
        if ind % (rows * cols) == 0:
            if ind == 0:
                if new_page:
                    new_slide = add_new_content_page(presentation, params)
                else:
                    new_slide = presentation.slides[-1]
            else:
                new_slide = add_new_content_page(presentation, params)
            page_ind += 1

        # grid position
        x_multiple = ind % cols
        y_multiple = ind // cols
        d_pos_x = pos_x + x_multiple * grid_w
        d_pos_y = pos_y + (y_multiple - page_ind * rows) * grid_h

        # grid chart data
        plot_textbox(new_slide.shapes, d_pos_x, d_pos_y, grid_w, grid_h, params)
        add_footnotes(new_slide, params.get("footnotes", []))


def add_chart_table_page(
    presentation,
    chart_data: Union[pd.DataFrame, list, ChartData, XyChartData],
    table_data: Union[pd.DataFrame, pd.Series],
    params: dict = {},
):
    new_slide = add_new_content_page(presentation, params)

    ## plot the bars of recent perforamnces
    # chart position and dimension
    chart_x = BODY_X
    chart_y = TITLE_Y + TITLE_HEIGHT
    chart_width = BODY_WIDTH
    chart_height = BODY_HEIGHT * params.get("chart", {}).get("height_pct", 1 / 3)

    print(params)
    # chart params
    # add the perf column chart
    if params.get("chart", {}).get("type", "line") == "column":
        chart_params = params.get("chart", {})
        #print(chart_params)
        plot_clustered_columns(
            new_slide,
            chart_data,
            chart_x,
            chart_y,
            chart_width,
            chart_height,
            chart_params,
        )
    else:
        c_data = line_chart_data(chart_data, params)

        chart_params = add_params(
            params.get("chart", {}),
            c_data["params"],
        )

        plot_chart(
            new_slide.shapes,
            c_data["data"],
            chart_x,
            chart_y,
            chart_width,
            chart_height,
            chart_params.get("type", "line"),
            chart_params,
        )

    ## plot the table of perforamnce overview
    # set table position and dimension
    table_x = BODY_X
    table_y = chart_y + chart_height + MARGIN_VER
    table_width = BODY_WIDTH
    table_height = PAGE_HEIGHT - table_y

    table_params = params.get("table", {"percentage": True})
    # add the table
    plot_table(
        new_slide, table_data, table_x, table_y, table_width, table_height, table_params
    )
    add_footnotes(new_slide, params.get("footnotes", []))


def delete_slides(prs, descriptions_slides, slide_to_keep):
    for (
        i
    ) in (
        descriptions_slides
    ):  # remmove slides 15-16-17 dans un ppt (14-15-16 selon numerotation python)
        if i != slide_to_keep:
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]


def ppt_to_pdf(
    source_filename,
    output_filename,
    # output_filename2,
    formatType=32,
):
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    powerpoint.Visible = 1

    if output_filename[-3:] != "pdf":
        output_filename = output_filename + ".pdf"
    deck = powerpoint.Presentations.Open(source_filename)
    deck.SaveAs(output_filename, formatType)  # formatType = 32 for ppt to pdf
    # deck.SaveAs(output_filename2, formatType)  # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()
