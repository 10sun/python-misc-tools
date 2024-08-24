'''
Author: J , jwsun1987@gmail.com
Date: 2024-02-09 18:00:08
LastEditors: J , jwsun1987@gmail.com
Description: 
Copyright: Copyright (c) 2024 by jwsun1987@gmail.com. All Rights Reserved.
'''


try:
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    import sys

    sys.path.append(
        r"\\merlin\lib_isg\\28.Alternative Data\Code\python-quant\libs\common\reporting\slides\python-pptx-master"
    )
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

from .lo_slides_style import *
#from common.reporting.style.lo_color import *

table_types = [
    "{2D5ABB26-0587-4C30-8999-92F81FD0307C}",
    "{3C2FFA5D-87B4-456A-9821-1D502468CF0F}",
    "{284E427A-3D55-4303-BF80-6455036E1DE7}",
    "{69C7853C-536D-4A76-A0AE-DD22124D55A5}",
    "{775DCB02-9BB8-47FD-8907-85C794F793BA}",
    "{35758FB7-9AC5-4552-8A53-C91805E547FA}",
    "{08FB837D-C827-4EFA-A057-4D05807E0F7C}",
    "{5940675A-B579-460E-94D1-54222C63F5DA}",
    "{D113A9D2-9D6B-4929-AA2D-F23B5EE8CBE7}",
    "{18603FDC-E32A-4AB5-989C-0864C3EAD2B8}",
    "{306799F8-075E-4A3A-A7F6-7FBC6576F1A4}",
    "{E269D01E-BC32-4049-B463-5C60D7B0CCD2}",
    "{327F97BB-C833-4FB7-BDE5-3F7075034690}",
    "{638B1855-1B75-4FBE-930C-398BA8C253C6}",
    "{9D7B26C5-4107-4FEC-AEDC-1716B250A1EF}",
    "{3B4B98B0-60AC-42C2-AFA5-B58CD77FA1E5}",
    "{0E3FDE45-AF77-4B5C-9715-49D594BDF05E}",
    "{C083E6E3-FA7D-4D7B-A595-EF9225AFEA82}",
    "{D27102A9-8310-4765-A935-A1911B00CA55}",
    "{5FD0F851-EC5A-4D38-B0AD-8093EC10F338}",
    "{68D230F3-CF80-4859-8CE7-A43EE81993B5}",
    "{7E9639D4-E3E2-4D34-9284-5A2195B3D0D7}",
    "{69012ECD-51FC-41F1-AA8D-1B2483CD663E}",
    "{72833802-FEF1-4C79-8D5D-14CF1EAF98D9}",
    "{F2DE63D5-997A-4646-A377-4702673A728D}",
    "{17292A2E-F333-43FB-9621-5CBBE7FDCDCB}",
    "{5A111915-BE36-4E01-A7E5-04B1672EAD32}",
    "{912C8C85-51F0-491E-9774-3900AFEF0FD7}",
    "{616DA210-FB5B-4158-B5E0-FEB733F419BA}",
    "{BC89EF96-8CEA-46FF-86C4-4CE0E7609802}",
    "{5DA37D80-6434-44D0-A028-1B22A696006F}",
    "{8799B23B-EC83-4686-B30A-512413B5E67A}",
    "{ED083AE6-46FA-4A59-8FB0-9F97EB10719F}",
    "{BDBED569-4797-4DF1-A0F4-6AAB3CD982D8}",
    "{E8B1032C-EA38-4F05-BA0D-38AFFFC7BED3}",
    "{793D81CF-94F2-401A-BA57-92F5A7B2D0C5}",
    "{B301B821-A1FF-4177-AEE7-76D212191A09}",
    "{9DCAF9ED-07DC-4A11-8D7F-57B35C25682E}",
    "{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}",
    "{1E171933-4619-4E11-9A3F-F7608DF75F80}",
    "{FABFCF23-3B69-468F-B69F-88F6DE6A72F2}",
    "{10A1B5D5-9B99-4C35-A422-299274C87663}",
    "{073A0DAA-6AF3-43AB-8588-CEC1D06C72B9}",
    "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}",
    "{21E4AEA4-8DFA-4A89-87EB-49C32662AFE0}",
    "{F5AB1C69-6EDB-4FF4-983F-18BD219EF322}",
    "{00A15C55-8517-42AA-B614-E9B94910E393}",
    "{7DF18680-E054-41AD-8BC1-D1AEF772440D}",
    "{93296810-A885-4BE3-A3E7-6D5BEEA58F35}",
    "{8EC20E35-A176-4012-BC5E-935CFFF8708E}",
    "{6E25E649-3F16-4E02-A733-19D2CDBF48F0}",
    "{85BE263C-DBD7-4A20-BB59-AAB30ACAA65A}",
    "{EB344D84-9AFB-497E-A393-DC336BA19D2E}",
    "{EB9631B5-78F2-41C9-869B-9F39066F8104}",
    "{74C1A8A3-306A-4EB7-A6B1-4F7E0EB9C5D6}",
    "{2A488322-F2BA-4B5B-9748-0D474271808F}",
    "{D7AC3CCA-C797-4891-BE02-D94E43425B78}",
    "{69CF1AB2-1976-4502-BF36-3FF5EA218861}",
    "{8A107856-5554-42FB-B03E-39F5DBC370BA}",
    "{0505E3EF-67EA-436B-97B2-0124C06EBD24}",
    "{C4B1156A-380E-4F78-BDF5-A606A8083BF9}",
    "{22838BEF-8BB2-4498-84A7-C5851F593DF1}",
    "{16D9F66E-5EB9-4882-86FB-DCBF35E3C3E4}",
    "{E8034E78-7F5D-4C2E-B375-FC64B27BC917}",
    "{125E5076-3810-47DD-B79F-674D7AD40C01}",
    "{37CE84F3-28C3-443E-9E96-99CF82512B78}",
    "{D03447BB-5D67-496B-8E87-E561075AD55C}",
    "{E929F9F4-4A8F-4326-A1B4-22849713DDAB}",
    "{8FD4443E-F989-4FC4-A0C8-D5A2AF1F390B}",
    "{AF606853-7671-496A-8E4F-DF71F8EC918B}",
    "{5202B0CA-FC54-4496-8BCA-5EF66A818D29}",
    "{0660B408-B3CF-4A94-85FC-2B1E0A45F4A2}",
    "{91EBBBCC-DAD2-459C-BE2E-F6DE35CF9A28}",
    "{46F890A9-2807-4EBB-B81D-B2AA78EC7F39}",
]


def set_table_style(shape, params={}):
    tbl = shape._element.graphic.graphicData.tbl
    tbl[0][-1].text = params.get(
        "style_id", "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
    )  # 3B4B98B0-60AC-42C2-AFA5-B58CD77FA1E5


def set_cell_margins(cell, params={}):
    if params.get("top", None) is not None:
        cell.margin_top = params.get("top")
    if params.get("bottom", None) is not None:
        cell.margin_bottom = params.get("bottom")
    if params.get("left", None) is not None:
        cell.margin_left = params.get("left")
    if params.get("right", None) is not None:
        cell.margin_right = params.get("right")


def set_table_content(table, df, params={}):
    n_header_lvl = df.columns.nlevels
    n_index_lvl = df.index.nlevels
    # fill in the df content
    for r in range(0, df.shape[0]):
        for c in range(0, df.shape[1]):
            table.cell(r + n_header_lvl, c + n_index_lvl).text = (
                str(df.iloc[r, c]) + "%"
                if params.get("percentage", False)
                else str(df.iloc[r, c])
            )
            table.cell(r + n_header_lvl, c + n_index_lvl).text_frame.paragraphs[
                0
            ].font.size = params.get("body_size", TABLE_BODY_FONT_SIZE)
            table.cell(r + n_header_lvl, c + n_index_lvl).text_frame.paragraphs[
                0
            ].alignment = PP_ALIGN.CENTER
            set_cell_margins(
                table.cell(r + n_header_lvl, c + n_index_lvl),
                params.get("cell_margins", {}),
            )

    # fill the index
    for r in range(0, df.shape[0]):
        for ind_col in range(0, n_index_lvl):
            table.cell(r + n_header_lvl, ind_col).text = (
                df.index[r] if n_index_lvl == 1 else df.index[r][ind_col]
            )
            table.cell(r + n_header_lvl, ind_col).text_frame.paragraphs[
                0
            ].font.size = params.get("index_size", TABLE_HEADER_FONT_SIZE)
            table.cell(r + n_header_lvl, ind_col).text_frame.paragraphs[
                0
            ].alignment = PP_ALIGN.LEFT
            table.cell(r + n_header_lvl, ind_col).fill.solid()
            table.cell(r + n_header_lvl, ind_col).fill.fore_color.rgb = params.get(
                "index_fill_color", LO_LIGHTGRAY
            )

    # fill the header
    for c in range(0, df.shape[1]):
        for header_row in range(0, n_header_lvl):
            table.cell(header_row, c + n_index_lvl).text = (
                df.columns[c] if n_header_lvl == 1 else df.columns[c][header_row]
            )
            table.cell(header_row, c + n_index_lvl).text_frame.paragraphs[
                0
            ].font.size = params.get("header_size", TABLE_HEADER_FONT_SIZE)
            table.cell(header_row, c + n_index_lvl).text_frame.paragraphs[
                0
            ].font.color.rgb = params.get("header_color", RGBColor(255, 255, 255))
            table.cell(header_row, c + n_index_lvl).text_frame.paragraphs[
                0
            ].alignment = PP_ALIGN.CENTER
            # set fill type to solid color first
            table.cell(header_row, c + n_index_lvl).fill.solid()
            # set foreground (fill) color to a specific RGB color
            table.cell(header_row, c + n_index_lvl).fill.fore_color.rgb = params.get(
                "header_fill_color", LO_BROWN
            )

    # merge the header if it is the same as its neighbor
    for header_lvl in range(0, n_header_lvl):

        header_col = 0
        next_col = header_col + 1

        while next_col < df.shape[1] + 1:
            curr_cell = table.cell(header_lvl, header_col)
            curr_text = curr_cell.text

            if curr_text == table.cell(header_lvl, next_col).text:
                if next_col == df.shape[1]:
                    merge_cells(
                        table,
                        header_lvl,
                        header_col,
                        header_lvl,
                        next_col,
                        {
                            **params,
                            **{
                                "header_size": TABLE_TITLE_FONT_SIZE,
                                "text_align": PP_ALIGN.CENTER,
                            },
                        },
                    )
            elif next_col > header_col + 1:
                merge_cells(
                    table,
                    header_lvl,
                    header_col,
                    header_lvl,
                    next_col - 1,
                    {
                        **params,
                        **{
                            "header_size": TABLE_TITLE_FONT_SIZE,
                            "text_align": PP_ALIGN.CENTER,
                        },
                    },
                )
                header_col = next_col
                # auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            else:
                header_col += 1
            next_col += 1

    # merge the index if it is empty

    for ind_lvl in range(0, n_index_lvl):
        header_row = n_index_lvl
        next_row = header_row + 1

        while next_row < df.shape[0] + 1:
            curr_cell = table.cell(header_row, ind_lvl)
            curr_text = curr_cell.text

            if curr_text == table.cell(next_row, ind_lvl).text:
                if next_row == df.shape[0]:
                    merge_cells(table, header_row, ind_lvl, next_row, ind_lvl, params)
            elif next_row > header_row + 1:
                merge_cells(table, header_row, ind_lvl, next_row - 1, ind_lvl, params)
                header_row = next_row
            else:
                header_row += 1
            next_row += 1

    # merge the top left corner if there is more than one level of header
    for row in range(n_header_lvl):
        for col in range(n_index_lvl):
            table.cell(row, col).fill.solid()
            table.cell(row, col).fill.fore_color.rgb = params.get(
                "header_fill_color", LO_BROWN
            )


def set_cell_dimension(cell, params={}):

    pass


def set_table_dimension(table, params={}):
    pass


def merge_cells(table, org_row, org_col, end_row, end_col, params: dict = {}):
    curr_cell = table.cell(org_row, org_col)
    curr_text = curr_cell.text
    curr_cell.merge(table.cell(end_row, end_col))
    curr_cell.text = curr_text
    curr_cell.text_frame.paragraphs[0].font.size = params.get(
        "header_size", TABLE_HEADER_FONT_SIZE
    )
    curr_cell.text_frame.paragraphs[0].alignment = params.get(
        "text_align", PP_ALIGN.LEFT
    )
    curr_cell.text_frame.paragraphs[0].font.color.rgb = params.get(
        "header_color", RGBColor(255, 255, 255)
    )


def color_code_table(table, params={}):
    from colour import Color

    colors = list(red.range_to(Color("green"), 10))
