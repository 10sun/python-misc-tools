'''
Author: J , jwsun1987@gmail.com
Date: 2024-02-09 17:59:04
LastEditors: J , jwsun1987@gmail.com
Description: 
Copyright: Copyright (c) 2024 by jwsun1987@gmail.com. All Rights Reserved.
'''

import numpy as np

try:
    import pptx
    from pptx import enum as pptx_enum
except ImportError:
    import sys

    sys.path.append(
        r"\\merlin\lib_isg\\28.Alternative Data\Code\python-quant\libs\common\reporting\slides\python-pptx-master"
    )
    import pptx
    from pptx import enum as pptx_enum

#from common.reporting.style.lo_color import *
#from common.reporting.style.lo_style import *

from .lo_slides_style import *


def translate_chart_type(chart_type_str: str) -> pptx_enum.chart.XL_CHART_TYPE:
    """[summary]

    Args:
        chart_type_str (str): [description]

    Returns:
        chart_enum.XL_CHART_TYPE: [description]
    """
    if chart_type_str.casefold() == "line".casefold():
        chart_type = pptx_enum.chart.XL_CHART_TYPE.LINE
    elif chart_type_str.casefold() == "column".casefold():
        chart_type = pptx_enum.chart.XL_CHART_TYPE.COLUMN_CLUSTERED
    elif chart_type_str.casefold() == "bar".casefold():
        chart_type = pptx_enum.chart.XL_CHART_TYPE.BAR_CLUSTERED
    elif chart_type_str.casefold() == "bubble".casefold():
        chart_type = pptx_enum.chart.XL_CHART_TYPE.BUBBLE
    elif chart_type_str.casefold() == "pie".casefold():
        chart_type = pptx_enum.chart.XL_CHART_TYPE.PIE
    elif chart_type_str.casefold() == "radar".casefold():
        chart_type = pptx_enum.chart.XL_CHART_TYPE.RADAR
    elif chart_type_str.casefold() == "surface".casefold():
        chart_type = pptx_enum.chart.XL_CHART_TYPE.SURFACE
    elif chart_type_str.casefold() == "scatter".casefold():
        chart_type = pptx_enum.chart.XL_CHART_TYPE.XY_SCATTER
    else:
        ValueError(chart_type_str + " not available...")
        return
    return chart_type


def set_chart_font(chart, params: dict):
    """set the font of the chart

    Args:
        chart ([type]): [description]
        params (dict): [description]
    """
    chart.font.name = params.get("font", {}).get("name", FONT_NAME)
    chart.font.size = params.get("font", {}).get("size", CHART_BODY_FONT_SIZE)
    chart.font.rgb = params.get("font", {}).get("color", CHART_FONT_COLOR)


def set_chart_title(chart, params: dict):
    """set the title of the chart

    Args:
        chart ([type]): [description]
        params (dict): [description]
    """
    chart.has_title = params.get("on", True)
    if chart.has_title:
        chart.chart_title.text_frame.text = params.get("text", "")
        chart.chart_title.text_frame.paragraphs[0].font.name = params.get(
            "font", FONT_NAME
        )
        chart.chart_title.text_frame.paragraphs[0].font.size = params.get(
            "font_size", CHART_HEADER_FONT_SIZE
        )
        chart.chart_title.text_frame.paragraphs[0].font.bold = params.get(
            "font_bold", True
        )


def set_chart_legend(chart, params: dict):
    """set the legend of the chart

    Args:
        chart ([type]): [description]
        params (dict): [description]
    """
    chart.has_legend = params.get("on", True)
    if chart.has_legend:
        chart.legend.position = params.get(
            "position", pptx_enum.chart.XL_LEGEND_POSITION.TOP
        )
        chart.legend.include_in_layout = params.get("in", False)
        chart.legend.font.name = params.get("font", FONT_NAME)
        chart.legend.font.size = params.get("font_size", CHART_BODY_FONT_SIZE)
        chart.legend.font.bold = params.get("font_bold", False)


def set_chart_data_labels(chart, params: dict):
    """set the data labels of the chart

    Args:
        chart ([type]): [description]
        params (dict): [description]
    """
    plot = chart.plots[0]
    plot.has_data_labels = params.get("on", False)
    if plot.has_data_labels:
        if params.get("how", "all") == "all":
            data_labels = plot.data_labels
            data_labels.position = params.get(
                "position", pptx_enum.chart.XL_DATA_LABEL_POSITION.OUTSIDE_END
            )
            data_labels.show_percentage = params.get("percentage", False)
            if data_labels.show_percentage:
                for series in chart.series:
                    for point, value in zip(series.points, series.values):
                        point.data_label.text_frame.text = str(
                            "{:.2%}".format(value / 100)
                        )
                # data_labels.number_format = params.get("number_format", '0"%"')
        elif params.get("how", "trough") == "trough":
            for series_idx, series in enumerate(chart.series):
                trough_idx = np.argmin(series.values)  # == min(series.values))
                data_label = series.points[trough_idx].data_label
                data_label.position = params.get(
                    "position", pptx_enum.chart.XL_DATA_LABEL_POSITION.BELOW
                )
                data_label.has_text_frame = True
                if params.get("percentage", False):
                    data_label.text_frame.paragraphs[0].text = (
                        str(round(min(series.values), 1)) + "%"
                    )
                else:
                    data_label.text_frame.paragraphs[0].text = str(
                        round(min(series.values), 1)
                    )
                data_label.text_frame.paragraphs[0].runs[0].font.name = params.get(
                    "font", FONT_NAME
                )
                data_label.text_frame.paragraphs[0].runs[0].font.size = params.get(
                    "font_size", CHART_HEADER_FONT_SIZE
                )
                data_label.text_frame.paragraphs[0].runs[0].font.bold = params.get(
                    "font_bold", True
                )
                data_label.text_frame.paragraphs[0].runs[
                    0
                ].font.color.rgb = COUNTRY_COLOR.get(
                    series.name, LO_COLOR_PALETTE[series_idx]
                )


def set_chart_axis(axis, params: dict):
    """set the axis of the chart

    Args:
        axis ([type]): [description]
        params (dict): [description]
    """
    # visibility
    axis.visible = params.get("visible", True)

    if isinstance(axis, pptx.chart.axis.DateAxis):
        axis.major_unit = params.get("major_unit", 12)
        axis.minor_unit = axis.major_unit / 2
        axis.major_time_unit = params.get(
            "major_time_unit", pptx_enum.chart.XL_TIME_UNIT.MONTHS
        )

    # grid lines
    axis.has_major_gridlines = params.get("major_gridlines", False)
    axis.has_minor_gridlines = params.get("minor_gridlines", False)

    # tick mark
    axis.major_tick_mark = params.get(
        "major_tick_mark", pptx_enum.chart.XL_TICK_MARK.INSIDE
    )
    axis.minor_tick_mark = params.get(
        "minor_tick_mark", pptx_enum.chart.XL_TICK_MARK.NONE
    )

    # axis scale
    if params.get("maximum_scale", np.nan) is not np.nan:
        axis.maximum_scale = params.get("maximum_scale", None)
    if params.get("minimum_scale", np.nan) is not np.nan:
        axis.minimum_scale = params.get("minimum_scale", None)

    # reverse order
    axis.reverse_order = params.get("reverse", False)

    # axis title
    if params.get("title", {}).get("on", False):
        # axis.axis_title.top = Cm(0)
        # axis.axis_title.left = Cm(0)
        axis.axis_title.text_frame.text = params.get("title", {}).get(
            "text", ValueError("No axis label")
        )
        axis.axis_title.text_frame.paragraphs[0].font.name = params.get(
            "title", {}
        ).get("font", FONT_NAME)
        axis.axis_title.text_frame.paragraphs[0].font.size = params.get(
            "title", {}
        ).get("font_size", CHART_TITLE_FONT_SIZE)
        axis.axis_title.text_frame.paragraphs[0].font.bold = params.get(
            "title", {}
        ).get("font_bold", True)
        axis.axis_title.text_frame.auto_size = (
            pptx_enum.text.MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        )


def set_chart_axis_ticks(axis, params: dict):
    """set the axis ticks of the chart

    Args:
        axis ([type]): [description]
        params (dict): [description]
    """
    axis.tick_label_position = params.get(
        "tick_label_position", pptx_enum.chart.XL_TICK_LABEL_POSITION.LOW
    )

    tick_labels = axis.tick_labels
    if params.get("percentage", False):
        # label.data_label.text_frame.text = str("{:.1%}".format(value*100))
        tick_labels.number_format = params.get("number_format", '0.0"%"')
    elif isinstance(axis, pptx.chart.axis.DateAxis):
        tick_labels.number_format = params.get("number_format", "mm.yyyy")
    tick_labels.font.bold = params.get("font_bold", False)
    tick_labels.font.size = params.get("font_size", CHART_BODY_FONT_SIZE)
    if params.get("offset", None) is not None and isinstance(
        axis, pptx.chart.axis.CategoryAxis
    ):
        tick_labels.offset = params.get("offset")


def set_chart_style(chart, chart_type: pptx_enum.chart.XL_CHART_TYPE, params: dict):
    """set the chart style

    Args:
        chart ([type]): [description]
        chart_type (chart_enum.XL_CHART_TYPE): [description]
        params (dict): [description]
    """
    if chart_type == pptx_enum.chart.XL_CHART_TYPE.LINE:
        for series in chart.series:
            series.format.line.width = params.get("line", {}).get(
                "width", CHART_LINE_WIDTH
            )
            if "mean".casefold() in series.name.casefold():
                series.format.line.dash_style = pptx_enum.dml.MSO_LINE.LONG_DASH
                series.format.line.color.rgb = LO_DARKGRAYBLUE
            elif "1std".casefold() in series.name.casefold():
                series.format.line.dash_style = pptx_enum.dml.MSO_LINE.ROUND_DOT
                series.format.line.color.rgb = LO_DARKGRAYBLUE
            elif "2std".casefold() in series.name.casefold():
                series.format.line.dash_style = pptx_enum.dml.MSO_LINE.ROUND_DOT
                series.format.line.color.rgb = LO_DARKGRAYBLUE
            elif series.name in COUNTRY_COLOR.keys() and params.get("line", {}).get(
                "country", True
            ):
                series.format.line.color.rgb = COUNTRY_COLOR.get(series.name)
                series.format.line.width = CHART_BOLD_LINE_WIDTH
            elif params.get("line", {}).get(series.name, False):
                series.format.line.dash_style = (
                    params.get("line", {})
                    .get(series.name, {})
                    .get("style", pptx_enum.dml.MSO_LINE.SOLID)
                )
                series.format.line.color.rgb = (
                    params.get("line", {}).get(series.name, {}).get("color", LO_NAVY)
                )
                series.format.line.width = (
                    params.get("line", {})
                    .get(series.name, {})
                    .get("width", CHART_LINE_WIDTH)
                )
    elif chart_type == pptx_enum.chart.XL_CHART_TYPE.COLUMN_CLUSTERED:
        palette = LO_COLOR_PALETTE[::-1]
        print(palette)
        for series in chart.series:
            if not palette:
                palette = LO_COLOR_PALETTE[::-1]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = params.get("bar", {}).get(
                "color", palette.pop()
            )
            series.invert_if_negative = params.get("bar", {}).get("invert", False)
